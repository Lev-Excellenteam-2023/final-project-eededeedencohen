import os
import collections.abc
from pptx import Presentation
import asyncio
from pptx.exc import PackageNotFoundError


def is_valid_pptx(pptx_file_path: str) -> bool:
    """
    @summary:
        Validates the file path and the file extension to be .pptx
    @param:
        str pptx_file_path: path to pptx file
    @return:
        bool: True if valid pptx file, False otherwise
    """
    if not os.path.exists(pptx_file_path):
        return False
    if not pptx_file_path.endswith('.pptx'):
        return False
    try:
        Presentation(pptx_file_path)
    except PackageNotFoundError:
        return False
    return True


async def get_list_of_content_from_pptx_file(pptx_file_path: str) -> list:
    """
    @summary:
        Receive a path to a pptx file and return a list of content from the pptx file.
        Each content is a list of sections from a slide.
        Example:
        get_list_of_content_from_pptx_file("some_file.pptx") ->
        [
        [ [section1 of slide1], [section2 of slide1], ... ,[sectionN of slide1] ],
        [ [section1 of slide2], [section2 of slide2], ... ,[sectionN of slide2] ], ... ,
        [ [section1 of slideN], [section2 of slideN], ... ,[sectionN of slideN] ]
        ]
        where section1 is the title of the slide and the first slide is the title of the pptx file.
    @param:
        str pptx_file_path: path to pptx file
    @return:
        list: list of content from pptx file
    @raise:
        PackageNotFoundError: If the pptx_file_path is not a valid path to a pptx file.
    """
    if not is_valid_pptx(pptx_file_path):
        raise ValueError(f"{pptx_file_path} is not a valid path to a pptx file.")

    loop = asyncio.get_event_loop()
    ppt = await loop.run_in_executor(None, Presentation, pptx_file_path)
    content_list = []
    for slide in ppt.slides:
        slide_text = [shape.text and shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        slide_text_separated_by_sections = [section for text in slide_text for section in text.split('\n')]
        clean_slide_text = [[text] for text in slide_text_separated_by_sections if text]
        if clean_slide_text:
            content_list.append(clean_slide_text)
    return content_list


async def main():
    pptx_file_path = "./files/asyncio-intro.pptx"
    content_list = await get_list_of_content_from_pptx_file(pptx_file_path)
    for content in content_list:
        print(content)
    print("=====================================")
    print(content_list)


if __name__ == "__main__":
    asyncio.run(main())

