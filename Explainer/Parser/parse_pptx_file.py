import io
import os
import collections.abc
from pptx import Presentation
import asyncio
from pptx.exc import PackageNotFoundError
from DB.schema import get_pptx_bytes


def is_valid_pptx(pptx_file_path: str) -> bool:
    """
    @summary:
        Validates the file path and the file extension to be .pptx
    @param:
        str pptx_file_path: path to pptx file
    @return:
        bool: True if valid pptx file, False otherwise
    """
    if not os.path.exists(pptx_file_path):
        return False
    if not pptx_file_path.endswith('.pptx'):
        return False
    try:
        Presentation(pptx_file_path)
    except PackageNotFoundError:
        return False
    return True


def extract_content_from_ppt(ppt) -> list:
    """
    Extracts content from the passed Presentation object.
    """
    content_list = []
    for slide in ppt.slides:
        slide_text = [shape.text and shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        slide_text_separated_by_sections = [section for text in slide_text for section in text.split('\n')]
        clean_slide_text = [[text] for text in slide_text_separated_by_sections if text]
        if clean_slide_text:
            content_list.append(clean_slide_text)
    return content_list


async def get_list_of_content_from_pptx_file(pptx_file_path: str) -> list:
    """
    @summary:
        Receive a path to a pptx file and return a list of content from the pptx file.
        Each content is a list of sections from a slide.
        Example:
        get_list_of_content_from_pptx_file("some_file.pptx") ->
        [
        [ [section1 of slide1], [section2 of slide1], ... ,[sectionN of slide1] ],
        [ [section1 of slide2], [section2 of slide2], ... ,[sectionN of slide2] ], ... ,
        [ [section1 of slideN], [section2 of slideN], ... ,[sectionN of slideN] ]
        ]
        where section1 is the title of the slide and the first slide is the title of the pptx file.
    @param:
        str pptx_file_path: path to pptx file
    @return:
        list: list of content from pptx file
    @raise:
        PackageNotFoundError: If the pptx_file_path is not a valid path to a pptx file.
    """
    if not is_valid_pptx(pptx_file_path):
        raise ValueError(f"{pptx_file_path} is not a valid path to a pptx file.")
    loop = asyncio.get_event_loop()
    ppt = await loop.run_in_executor(None, Presentation, pptx_file_path)
    return extract_content_from_ppt(ppt)


async def get_list_of_content_from_pptx_bytes(bytes_pptx_file: bytes) -> list:
    """
    @summary:
        Receive a bytes object of a pptx file and return a list of content.
        Each content is a list of sections from a slide.
    @param:
        bytes bytes_pptx_file: bytes object of a pptx file
    @return:
        list: list of content from pptx file
    @example:
        get_list_of_content_from_pptx_bytes(bytes_pptx_file) ->
        [
        [ [section1 of slide1], [section2 of slide1], ... ,[sectionN of slide1] ],
        [ [section1 of slide2], [section2 of slide2], ... ,[sectionN of slide2] ], ... ,
        [ [section1 of slideN], [section2 of slideN], ... ,[sectionN of slideN] ]
        ]
        where section1 is the title of the slide and the first slide is the title of the pptx file.
    """
    file_like_object = io.BytesIO(bytes_pptx_file)
    loop = asyncio.get_event_loop()
    pptx = await loop.run_in_executor(None, Presentation, file_like_object)
    return extract_content_from_ppt(pptx)


